import spacy
from pptx import Presentation
from django.core.files.base import ContentFile
import io

# Load the English language model
nlp = spacy.load('en_core_web_sm')

def extract_main_points(text):
    doc = nlp(text)
    sentences = list(doc.sents)
    main_points = [sent.text for sent in sentences if len(sent.text.split()) > 5]
    return main_points

def create_ppt(main_points):
    prs = Presentation()

    for point in main_points:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use slide layout index as needed
        slide.shapes.title.text = "Main Point"
        slide.placeholders[1].text = point

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def create_powerPoint(text):
    main_points = extract_main_points(text)
    ppt_io = create_ppt(main_points)
    return ContentFile(ppt_io.read(), 'output.pptx')
